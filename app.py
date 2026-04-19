import streamlit as st
import pandas as pd
import copy
import io
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ── 페이지 설정 ──────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="PPT 자동 생성기",
    page_icon="📊",
    layout="centered",
)

st.title("📊 PPT 자동 생성기")
st.caption("엑셀 데이터를 PPT 양식에 자동으로 입력합니다.")

# ── 사용 방법 안내 ────────────────────────────────────────────────────────────
with st.expander("📖 사용 방법", expanded=False):
    st.markdown("""
    **엑셀 파일 형식**

    | 직책 | 직책 폰트 | 직책 글씨크기 | 이름 | 이름 폰트 | 이름 글씨크기 |
    |------|----------|------------|------|----------|------------|
    | 대표이사 | 맑은 고딕 | 24 | 홍길동 | 맑은 고딕 | 36 |
    | 부장 | Arial | 20 | 김철수 | Arial | 28 |

    - 1행은 **헤더(컬럼명)** 입니다.
    - 2행부터 데이터가 시작되며, 각 행이 PPT 슬라이드 한 장에 해당합니다.

    **PPT 양식 파일**
    - 글상자 안에 `{직책}`, `{이름}` 플레이스홀더가 있어야 합니다.
    - 이미지는 모든 슬라이드에 동일하게 반복됩니다.
    - 엑셀 행 수만큼 슬라이드가 생성됩니다.
    """)

# ── 파일 업로드 ──────────────────────────────────────────────────────────────
col1, col2 = st.columns(2)

with col1:
    st.subheader("1️⃣ 엑셀 파일")
    excel_file = st.file_uploader(
        "엑셀 파일을 업로드하세요",
        type=["xlsx", "xls"],
        key="excel",
        help="직책, 직책 폰트, 직책 글씨크기, 이름, 이름 폰트, 이름 글씨크기 컬럼이 필요합니다."
    )

with col2:
    st.subheader("2️⃣ PPT 양식 파일")
    ppt_file = st.file_uploader(
        "PPT 파일을 업로드하세요",
        type=["pptx"],
        key="ppt",
        help="{직책}, {이름} 플레이스홀더가 글상자에 포함되어 있어야 합니다."
    )

# ── 엑셀 미리보기 ────────────────────────────────────────────────────────────
df = None
if excel_file:
    try:
        df = pd.read_excel(excel_file)

        # 컬럼명 정규화 (공백 제거)
        df.columns = df.columns.str.strip()

        required_cols = ["직책", "직책 폰트", "직책 글씨크기", "이름", "이름 폰트", "이름 글씨크기"]
        missing = [c for c in required_cols if c not in df.columns]

        if missing:
            st.error(f"❌ 엑셀 파일에 다음 컬럼이 없습니다: {', '.join(missing)}")
            df = None
        else:
            st.success(f"✅ 엑셀 파일 로드 완료 — {len(df)}개 행 (슬라이드 {len(df)}장 생성 예정)")
            st.dataframe(df, use_container_width=True)

    except Exception as e:
        st.error(f"❌ 엑셀 파일을 읽는 중 오류 발생: {e}")

# ── PPT 생성 핵심 함수 ────────────────────────────────────────────────────────
PLACEHOLDER_PATTERN = re.compile(r"\{(직책|이름)\}")

def get_font_name(val):
    """None / NaN 처리"""
    if val is None or (isinstance(val, float) and pd.isna(val)):
        return None
    return str(val).strip()

def get_font_size(val):
    """None / NaN 처리 후 Pt로 변환"""
    try:
        return Pt(float(val))
    except Exception:
        return None

# XML 네임스페이스
_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

def _set_run_text_xml(run, text):
    """
    <a:t> 태그의 텍스트만 교체한다.
    run.text = "..." 은 내부적으로 <a:r> 전체를 재구성해
    <a:rPr>(폰트 정보)를 날릴 수 있으므로, XML을 직접 조작한다.
    """
    t_elem = run._r.find(f"{{{_A}}}t")
    if t_elem is not None:
        t_elem.text = text
    else:
        # <a:t>가 없는 경우 새로 생성
        import lxml.etree as etree
        t_elem = etree.SubElement(run._r, f"{{{_A}}}t")
        t_elem.text = text

def _set_run_font_xml(run, font_name, font_size):
    """
    <a:rPr> 속성을 직접 수정해 폰트명·크기를 적용한다.
    없으면 새로 생성해서 <a:r>의 첫 번째 자식으로 삽입한다.
    font_size는 python-pptx Pt() 값(EMU 단위)이며,
    OOXML sz 속성은 '포인트 × 100' 정수값이다.
    """
    import lxml.etree as etree

    rPr = run._r.find(f"{{{_A}}}rPr")
    if rPr is None:
        rPr = etree.Element(f"{{{_A}}}rPr")
        run._r.insert(0, rPr)   # <a:rPr>은 <a:r>의 첫 번째 자식이어야 함

    if font_name:
        # <a:latin typeface="..."/> 처리
        latin = rPr.find(f"{{{_A}}}latin")
        if latin is None:
            latin = etree.SubElement(rPr, f"{{{_A}}}latin")
        latin.set("typeface", font_name)

    if font_size:
        # sz 속성: 포인트의 100배 정수 (예: 24pt → "2400")
        sz_val = str(int(round(font_size.pt * 100)))
        rPr.set("sz", sz_val)

def merge_and_replace_para(para, mapping):
    """
    paragraph 안의 run들을 하나로 합친 뒤 플레이스홀더를 치환한다.
    - run.text = ... 대신 XML <a:t> 직접 교체 → <a:rPr> 보존
    - 폰트도 <a:rPr> XML 직접 수정 → 덮어쓰기 없음
    """
    if not para.runs:
        return False

    # 1) paragraph 전체 텍스트 합산
    combined_text = "".join(r.text for r in para.runs)
    if not PLACEHOLDER_PATTERN.search(combined_text):
        return False

    # 2) 어떤 플레이스홀더가 있는지 파악
    found_keys = list(dict.fromkeys(PLACEHOLDER_PATTERN.findall(combined_text)))  # 순서 유지·중복 제거

    # 3) 첫 번째 run에 합산 텍스트를 XML 레벨로 넣고, 나머지 run은 비움
    first_run = para.runs[0]
    _set_run_text_xml(first_run, combined_text)
    for r in para.runs[1:]:
        _set_run_text_xml(r, "")

    # 4) 텍스트 치환 (XML 레벨)
    new_text = combined_text
    for key in found_keys:
        info = mapping.get(key)
        if info is None:
            continue
        new_text = new_text.replace(f"{{{key}}}", info["value"])
    _set_run_text_xml(first_run, new_text)

    # 5) 폰트 적용 — XML <a:rPr> 직접 수정 (텍스트 교체와 완전히 분리)
    for key in found_keys:
        info = mapping.get(key)
        if info:
            _set_run_font_xml(first_run, info.get("font_name"), info.get("font_size"))

    return True

def fill_slide(slide, mapping):
    """
    슬라이드 내 모든 텍스트프레임의 플레이스홀더를 치환한다.
    - shape이 몇 개든 상관없이 전부 순회
    - 한 shape 안에 동일 플레이스홀더가 여러 paragraph에 있어도 모두 처리
    - run이 쪼개져 있어도 paragraph 단위로 병합 후 치환
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # shape 전체 텍스트에 플레이스홀더가 있는지 빠르게 확인
        shape_full_text = "".join(
            r.text
            for para in shape.text_frame.paragraphs
            for r in para.runs
        )
        if not PLACEHOLDER_PATTERN.search(shape_full_text):
            continue

        # paragraph 단위로 처리 (한 shape 안에 여러 paragraph가 있어도 모두 처리)
        for para in shape.text_frame.paragraphs:
            merge_and_replace_para(para, mapping)

def generate_ppt(prs_template: Presentation, df: pd.DataFrame) -> bytes:
    """
    템플릿 첫 슬라이드를 기준으로 슬라이드를 복제해
    각 행의 데이터를 채워 넣는다.
    """
    # 새 Presentation 객체 (슬라이드 없이 시작)
    from pptx.util import Emu
    import lxml.etree as etree

    # 슬라이드 크기를 템플릿과 동일하게
    prs_out = Presentation()
    prs_out.slide_width = prs_template.slide_width
    prs_out.slide_height = prs_template.slide_height

    template_slide = prs_template.slides[0]

    for _, row in df.iterrows():
        mapping = {
            "직책": {
                "value": str(row["직책"]).strip() if pd.notna(row["직책"]) else "",
                "font_name": get_font_name(row.get("직책 폰트")),
                "font_size": get_font_size(row.get("직책 글씨크기")),
            },
            "이름": {
                "value": str(row["이름"]).strip() if pd.notna(row["이름"]) else "",
                "font_name": get_font_name(row.get("이름 폰트")),
                "font_size": get_font_size(row.get("이름 글씨크기")),
            },
        }

        # 슬라이드 레이아웃(blank) 추가
        blank_layout = prs_out.slide_layouts[6]  # blank
        new_slide = prs_out.slides.add_slide(blank_layout)

        # ── spTree 전체를 한 번에 deepcopy → 완전히 독립된 XML 트리 ──────
        # 자식 단위 deepcopy는 lxml 내부 참조가 공유되어
        # 두 번째 글상자 이후 치환이 첫 번째 슬라이드에도 반영되는 버그 발생.
        # spTree 전체를 한 번에 deepcopy해야 슬라이드마다 완전히 격리된다.
        cloned_sp_tree = copy.deepcopy(template_slide._element.spTree)

        dst_sp_tree = new_slide._element.spTree
        for child in list(dst_sp_tree):   # 기존 요소 전체 제거
            dst_sp_tree.remove(child)
        for child in list(cloned_sp_tree):  # 독립 복사본 삽입
            dst_sp_tree.append(child)

        # 이미지 관계(rId) 복사: 템플릿의 미디어를 새 슬라이드에 등록
        for rel in template_slide.part.rels.values():
            if "image" in rel.reltype:
                img_part = rel.target_part
                new_slide.part.relate_to(img_part, rel.reltype)

        # 텍스트 치환
        fill_slide(new_slide, mapping)

    # 바이트로 저장
    buf = io.BytesIO()
    prs_out.save(buf)
    buf.seek(0)
    return buf.read()

# ── 생성 버튼 ────────────────────────────────────────────────────────────────
st.divider()

if df is not None and ppt_file is not None:
    if st.button("🚀 PPT 생성하기", type="primary", use_container_width=True):
        with st.spinner("PPT 생성 중..."):
            try:
                ppt_bytes = ppt_file.read()
                prs_template = Presentation(io.BytesIO(ppt_bytes))

                result_bytes = generate_ppt(prs_template, df)

                st.success(f"✅ PPT 생성 완료! 슬라이드 {len(df)}장이 만들어졌습니다.")
                st.download_button(
                    label="⬇️ PPT 다운로드",
                    data=result_bytes,
                    file_name="output.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
            except Exception as e:
                st.error(f"❌ PPT 생성 중 오류 발생: {e}")
                st.exception(e)
else:
    st.info("📂 엑셀 파일과 PPT 양식 파일을 모두 업로드하면 생성 버튼이 활성화됩니다.")
